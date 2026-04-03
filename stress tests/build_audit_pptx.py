from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

DARK    = RGBColor(0x1a, 0x1a, 0x2e)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
GREEN   = RGBColor(0x1a, 0x7a, 0x3c)
GREEN_L = RGBColor(0xe8, 0xf5, 0xed)
GOLD    = RGBColor(0xC9, 0xA4, 0x3C)
AMBER   = RGBColor(0x92, 0x40, 0x0e)
AMBER_L = RGBColor(0xff, 0xfb, 0xeb)
RED     = RGBColor(0xb9, 0x1c, 0x1c)
MID     = RGBColor(0x33, 0x33, 0x33)
LIGHT   = RGBColor(0x66, 0x66, 0x66)
RULE    = RGBColor(0xdd, 0xdd, 0xdd)
TBL_H   = RGBColor(0x2d, 0x2d, 0x44)
TBL_A   = RGBColor(0xf7, 0xf7, 0xf9)
BLUE    = RGBColor(0x1e, 0x40, 0xaf)

blank = prs.slide_layouts[6]

def add_slide():
    return prs.slides.add_slide(blank)

def box(slide, x, y, w, h, bg=None, border=None):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.line.fill.background()
    if bg:
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg
    else:
        shape.fill.background()
    if border:
        shape.line.color.rgb = border
        shape.line.width = Pt(0.75)
    else:
        shape.line.fill.background()
    return shape

def txt(slide, text, x, y, w, h, size=11, bold=False, color=MID, align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.italic = italic
    return tb

def banner_box(slide, text, x, y, w, h, bg=GREEN_L, fg=GREEN, size=12):
    box(slide, x, y, w, h, bg=bg, border=fg)
    txt(slide, text, x+0.1, y+0.05, w-0.2, h-0.1, size=size, bold=True, color=fg, align=PP_ALIGN.CENTER)

def header_bar(slide, title, subtitle=''):
    box(slide, 0, 0, 13.33, 1.0, bg=DARK)
    txt(slide, title, 0.3, 0.05, 12, 0.55, size=22, bold=True, color=WHITE)
    if subtitle:
        txt(slide, subtitle, 0.3, 0.6, 12, 0.35, size=11, color=RGBColor(0xaa, 0xaa, 0xcc))

def table(slide, headers, rows, x, y, w, col_widths=None, row_h=0.3):
    if col_widths is None:
        col_widths = [w / len(headers)] * len(headers)
    cx = x
    for i, (h, cw) in enumerate(zip(headers, col_widths)):
        box(slide, cx, y, cw - 0.01, row_h, bg=TBL_H, border=RULE)
        txt(slide, h, cx + 0.06, y + 0.04, cw - 0.12, row_h - 0.08, size=8, bold=True, color=WHITE)
        cx += cw
    for ri, row in enumerate(rows):
        cy = y + row_h * (ri + 1)
        cx = x
        bg = WHITE if ri % 2 == 0 else TBL_A
        for i, (cell, cw) in enumerate(zip(row, col_widths)):
            box(slide, cx, cy, cw - 0.01, row_h, bg=bg, border=RULE)
            c = GREEN if cell in ('PASS', 'IDENTICAL', 'BASELINE', '0', 'YES') else (RED if cell in ('FAIL',) else MID)
            b = cell in ('PASS', 'IDENTICAL', 'BASELINE', 'TOTAL')
            txt(slide, str(cell), cx + 0.06, cy + 0.04, cw - 0.12, row_h - 0.08, size=8, bold=b, color=c)
            cx += cw

def check_list(slide, items, x, y, w, size=10):
    for i, item in enumerate(items):
        txt(slide, '✓  ' + item, x, y + i * 0.32, w, 0.3, size=size, color=GREEN)

# ── SLIDE 1 — TITLE ──────────────────────────────────────
sl = add_slide()
box(sl, 0, 0, 13.33, 7.5, bg=DARK)
txt(sl, 'Star Chart Converter', 0.5, 1.2, 12.33, 1.2, size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(sl, 'Complete Audit & Verification Record', 0.5, 2.5, 12.33, 0.7, size=22, color=RGBColor(0xaa, 0xaa, 0xcc), align=PP_ALIGN.CENTER)
txt(sl, 'Tool: Rise Athletics Star Chart Converter', 0.5, 3.4, 12.33, 0.4, size=13, color=RGBColor(0x88, 0x88, 0xaa), align=PP_ALIGN.CENTER)
txt(sl, 'Date: April 2, 2026     |     Audited by: Jayme Gibson + Claude (Anthropic)', 0.5, 3.85, 12.33, 0.4, size=12, color=RGBColor(0x88, 0x88, 0xaa), align=PP_ALIGN.CENTER)
banner_box(sl, 'FINAL RESULT: PASSED — All tests cleared. Zero errors across all audits.', 2.0, 5.0, 9.33, 0.6, bg=GREEN_L, fg=GREEN, size=14)

# ── SLIDE 2 — OVERVIEW ───────────────────────────────────
sl = add_slide()
header_bar(sl, 'Overview', 'Summary of all stress tests and audits performed')
table(sl,
    ['Test', 'Scope', 'Data Type', 'Gyms', 'Result'],
    [['Stress Test 1', 'Level 1, 6 classes', 'All scores filled (100%)', 'All 10', 'PASS'],
     ['Stress Test 1B', 'Level 1, 6 classes', 'Mixed scores + mastery', 'All 10', 'PASS'],
     ['Stress Test 3', 'All 7 programs, 47 pages', 'Real eval data, 216 students', 'CRR', '0 mismatches'],
     ['10-Gym Cross-Check', 'All 7 programs, 47 pages', 'Real eval data', 'All 10', 'IDENTICAL'],
    ],
    x=0.4, y=1.15, w=12.5,
    col_widths=[2.2, 2.4, 3.2, 1.5, 1.8],
    row_h=0.38
)
banner_box(sl, 'Grand Total: 11,132+ individual data points verified. Zero errors.', 1.5, 3.8, 10.33, 0.55, bg=GREEN_L, fg=GREEN, size=13)
txt(sl, 'Score Rendering Rules', 0.4, 4.55, 6, 0.35, size=13, bold=True, color=MID)
table(sl,
    ['Spreadsheet Value', 'Renders As', 'Meaning'],
    [['Blank cell (criteria row)', 'Empty grey circle', 'Not Yet'],
     ['1 in a criteria row', 'Filled colored star', 'Earned'],
     ['1 in the mastery / 3x in a row row', 'Gold star', 'Skill Complete'],
     ['Blank in mastery row', 'Empty grey circle', 'Not Yet / Not Mastered'],
    ],
    x=0.4, y=4.95, w=12.5,
    col_widths=[3.5, 3.0, 3.5],
    row_h=0.32
)
txt(sl, '* Student names: First Name + Last Initial only  |  * A/B splits render as separate pages  |  * Class names normalized from iClassPro format',
    0.4, 7.05, 12.5, 0.35, size=8, color=LIGHT)

# ── SLIDE 3 — STRESS TEST 1 ──────────────────────────────
sl = add_slide()
header_bar(sl, 'Stress Test 1 — Level 1, All Criteria Filled')
txt(sl, 'Purpose: Verify that when all criteria are earned (no blanks, no mastery), every star renders correctly, gym branding is correct, class names parse correctly, and A/B split logic works.',
    0.4, 1.1, 12.5, 0.45, size=10, color=MID)
txt(sl, 'Spreadsheet: Stress Test 1 - Level1 Full Scores No Signoff.xlsx  |  Score Pattern: All criteria = 1, all mastery = blank',
    0.4, 1.55, 12.5, 0.3, size=9, color=LIGHT, italic=True)
table(sl,
    ['Class', 'Students', 'Score Pattern'],
    [['Level 1  |  Monday  |  3:30pm', '8 (test names)', 'All criteria = 1, all mastery = blank'],
     ['Level 1  |  Monday  |  4:30pm', '8 (real names)', 'All criteria = 1, all mastery = blank'],
     ['Level 1  |  Monday  |  5:30pm', '8 (real names)', 'All criteria = 1, all mastery = blank'],
     ['Level 1  |  Monday  |  6:30pm', '6 (real names)', 'All criteria = 1, all mastery = blank'],
     ['Level 1 A  |  Monday  |  3:30pm', '8 (real names)', 'All criteria = 1, all mastery = blank'],
     ['Level 1 B  |  Monday  |  3:30pm', '7 (real names)', 'All criteria = 1, all mastery = blank'],
    ],
    x=0.4, y=1.9, w=12.5, col_widths=[4.5, 2.0, 6.0], row_h=0.33
)
txt(sl, 'Total: 1,980 earned (1s) + 720 blanks  |  Output: 10 gyms x 6 pages = 60 PDFs', 0.4, 4.05, 12.5, 0.3, size=9, color=LIGHT)
check_list(sl, [
    'All 10 gym names rendered correctly in headers',
    'All 6 class names parsed and formatted correctly',
    'A/B split at 3:30pm rendered as two separate pages',
    'All filled stars appear in correct gym color for each gym',
    'Mastery row shows empty circles where mastery = blank',
    'Student names display as first name + last initial only',
    'All 60 PDFs structurally identical except gym-specific branding',
], x=0.4, y=4.4, w=8.0, size=10)
banner_box(sl, 'STRESS TEST 1 RESULT: PASSED — 60 PDFs verified across all 10 gyms', 0.4, 6.85, 12.5, 0.5, size=12)

# ── SLIDE 4 — STRESS TEST 1B ─────────────────────────────
sl = add_slide()
header_bar(sl, 'Stress Test 1B — Level 1, Mixed Scores + Mastery')
txt(sl, 'Purpose: Verify that partial scores AND mastery stars both render correctly in the same PDF — including within the same class, side by side.',
    0.4, 1.1, 12.5, 0.4, size=10, color=MID)
txt(sl, 'Spreadsheet: Stress Test 1B  |  Same 6 classes and students as Stress Test 1',
    0.4, 1.5, 12.5, 0.3, size=9, color=LIGHT, italic=True)
txt(sl, 'Score Pattern — Two Groups in Every Class', 0.4, 1.9, 12.5, 0.35, size=12, bold=True, color=MID)

box(sl, 0.4, 2.3, 6.0, 1.6, bg=TBL_A, border=RULE)
txt(sl, 'Students 1-4  (Partial Scores)', 0.5, 2.35, 5.8, 0.3, size=10, bold=True, color=MID)
for i, line in enumerate(['Criteria 1 = 1  (earned)', 'Criteria 2 = 1  (earned)', 'Criteria 3 = blank  (not yet)', 'Mastery row = blank', 'PDF:  Star  Star  Circle  Circle']):
    c = BLUE if 'PDF' in line else MID
    b = 'PDF' in line
    txt(sl, line, 0.55, 2.65 + i * 0.22, 5.5, 0.22, size=9, color=c, bold=b)

box(sl, 6.7, 2.3, 6.0, 1.6, bg=GREEN_L, border=GREEN)
txt(sl, 'Students 5-8  (Full Scores + Mastery)', 6.8, 2.35, 5.8, 0.3, size=10, bold=True, color=GREEN)
for i, line in enumerate(['Criteria 1 = 1  (earned)', 'Criteria 2 = 1  (earned)', 'Criteria 3 = 1  (earned)', 'Mastery row = 1', 'PDF:  Star  Star  Star  GOLD Star']):
    c = GOLD if 'PDF' in line else MID
    b = 'PDF' in line
    txt(sl, line, 6.85, 2.65 + i * 0.22, 5.5, 0.22, size=9, color=c, bold=b)

txt(sl, 'Total: 1,890 earned (1s) + 810 blanks  |  Output: 10 gyms x 6 pages = 60 PDFs', 0.4, 4.05, 12.5, 0.3, size=9, color=LIGHT)
check_list(sl, [
    'Partial scores: C1 and C2 show as filled stars; C3 and mastery show as circles',
    'Full scores + mastery: all three criteria show as filled stars; mastery shows as gold star',
    'Gold stars are visually distinct from regular colored stars',
    'Mixed scoring within the same class renders correctly — different students, different patterns, same page',
    'All 10 gyms consistent. All 60 PDFs verified.',
], x=0.4, y=4.4, w=12.0, size=10)
banner_box(sl, 'STRESS TEST 1B RESULT: PASSED — Mixed scores and mastery stars verified across all 10 gyms', 0.4, 6.85, 12.5, 0.5, size=12)

# ── SLIDE 5 — STRESS TEST 3 ──────────────────────────────
sl = add_slide()
header_bar(sl, 'Stress Test 3 — Real Eval Data, All Programs, Full Visual Audit')
txt(sl, 'Purpose: The most rigorous test. Real iClassPro eval data from an actual eval date. Cell-by-cell visual audit of every star and circle in the output PDF vs. the source spreadsheet.',
    0.4, 1.1, 12.5, 0.4, size=10, color=MID)
txt(sl, 'Source: Real CRR iClassPro export, April 2, 2026  |  Output: CRR_Evals_04022026.pdf (47 pages)  |  Audit: 47 PNG images compared cell-by-cell. No cells skipped.',
    0.4, 1.5, 12.5, 0.3, size=9, color=LIGHT, italic=True)
table(sl,
    ['Program', 'Classes', 'Pages', 'Students', 'Skills/Student', 'Data Points', 'Mismatches'],
    [['Preschool', '5', '7 (with splits)', '26', '44', '1,144', '0'],
     ['Junior A', '4', '6 (with splits)', '28', '48', '1,344', '0'],
     ['Junior B', '4', '6 (with splits)', '23', '48', '1,104', '0'],
     ['Advanced Junior', '2', '2', '5', '58', '290', '0'],
     ['Level 1', '5', '7 (with splits)', '47', '58', '2,726', '0'],
     ['Level 2', '5', '9 (with splits)', '40', '52', '2,080', '0'],
     ['Level 3', '5', '10 (with splits)', '47', '52', '2,444', '0'],
     ['TOTAL', '30', '47 pages', '216', '—', '11,132', '0'],
    ],
    x=0.4, y=1.85, w=12.5,
    col_widths=[1.8, 1.1, 1.8, 1.2, 1.6, 1.8, 1.6],
    row_h=0.36
)
banner_box(sl, 'STRESS TEST 3 RESULT: 11,132 data points audited. ZERO errors. Every PDF matches its source spreadsheet exactly.', 0.4, 6.85, 12.5, 0.5, size=12)

# ── SLIDE 6 — 10-GYM CROSS-VALIDATION ───────────────────
sl = add_slide()
header_bar(sl, '10-Gym Cross-Validation')
txt(sl, 'Purpose: Confirm that the same eval data processed through all 10 gyms produces structurally identical output with only gym-specific branding differing.',
    0.4, 1.1, 12.5, 0.4, size=10, color=MID)
txt(sl, 'Method: Python/PyMuPDF extracted all text from all 47 pages of each gym PDF. Gym name headers stripped. Content compared page-by-page against CRR baseline.',
    0.4, 1.5, 12.5, 0.3, size=9, color=LIGHT, italic=True)
table(sl,
    ['Code', 'Gym Name', 'Pages', 'Name Correct', 'vs. Baseline'],
    [['CCP', 'Capital Gymnastics — Cedar Park', '47', 'YES', 'IDENTICAL'],
     ['CPF', 'Capital Gymnastics — Pflugerville', '47', 'YES', 'IDENTICAL'],
     ['CRR', 'Capital Gymnastics — Round Rock', '47', 'YES', 'BASELINE'],
     ['EST', 'Estrella Gymnastics', '47', 'YES', 'IDENTICAL'],
     ['HGA', 'Houston Gymnastics Academy', '47', 'YES', 'IDENTICAL'],
     ['OAS', 'Oasis Gymnastics', '47', 'YES', 'IDENTICAL'],
     ['RBA', 'Rowland Ballard — Atascocita', '47', 'YES', 'IDENTICAL'],
     ['RBK', 'Rowland Ballard — Kingwood', '47', 'YES', 'IDENTICAL'],
     ['SGT', 'Scottsdale Gymnastics', '47', 'YES', 'IDENTICAL'],
     ['TIG', 'Tigar Gymnastics', '47', 'YES', 'IDENTICAL'],
    ],
    x=1.5, y=1.9, w=10.0,
    col_widths=[1.0, 3.5, 1.0, 1.5, 2.0],
    row_h=0.35
)
txt(sl, '* EST (Estrella): Navy + Grey only — no red. Special branding handled correctly.', 0.4, 6.3, 12.5, 0.3, size=9, color=LIGHT)
txt(sl, '* Branding colors not captured by text extraction — star colors visually confirmed during Stress Tests 1 and 1B.', 0.4, 6.55, 12.5, 0.3, size=9, color=LIGHT)
banner_box(sl, '10-GYM RESULT: All 10 gyms — 47 pages each — identical structure. Only gym name header differs.', 0.4, 6.88, 12.5, 0.5, size=12)

# ── SLIDE 7 — PROGRAM NAME REFERENCE ────────────────────
sl = add_slide()
header_bar(sl, 'Program Name Reference — What the Converter Accepts', 'Internal reference. Case-insensitive, substring matching. Pipes stripped by iClassPro on export.')

txt(sl, 'Accepted Keywords', 0.4, 1.1, 8.0, 0.35, size=12, bold=True, color=MID)
table(sl,
    ['Keyword(s) in iClassPro class name', 'Resolves to'],
    [['PRESCHOOL / PRE-SCHOOL / PRE SCHOOL', 'Preschool'],
     ['JUNIOR / JR', 'Junior'],
     ['ADVANCED JUNIOR / ADV JUNIOR / ADV JR / ADVJR', 'Advanced Junior'],
     ['LEVEL 1 / LVL 1 / LEV1 / L1 / GIRLS LEVEL 1 / GIRLS L1 / GL1 / GIRLS 1', 'Level 1'],
     ['LEVEL 2 / LVL 2 / LEV2 / L2 / GIRLS LEVEL 2 / GIRLS L2 / GL2 / GIRLS 2', 'Level 2'],
     ['LEVEL 3 / LVL 3 / LEV3 / L3 / GIRLS LEVEL 3 / GIRLS L3 / GL3 / GIRLS 3', 'Level 3'],
    ],
    x=0.4, y=1.5, w=8.0, col_widths=[5.8, 2.2], row_h=0.35
)

txt(sl, 'What Gets Silently Skipped', 8.7, 1.1, 4.4, 0.35, size=12, bold=True, color=RED)
skipped = ['Recreational / Rec', 'Tots / Tiny Tots', 'Youth / Beginner',
           'Tumbling / Tumble', 'Boys / MAG / Mens',
           'Teen / Adult', 'Cheer / Acro / Dance / Ninja',
           'Level 4+', 'Level I / II / III (roman numerals)',
           'Ages 3-5 / coach name only']
for i, item in enumerate(skipped):
    txt(sl, '✗  ' + item, 8.7, 1.5 + i * 0.35, 4.3, 0.32, size=9, color=RED)

txt(sl, 'Known Risks', 0.4, 4.0, 12.5, 0.35, size=12, bold=True, color=AMBER)
box(sl, 0.4, 4.4, 12.5, 0.6, bg=AMBER_L, border=AMBER)
txt(sl, 'Silent wrong match — more dangerous than a skip. "Junior 2" contains "junior" so it maps to Junior program, not Level 2. PDF generates with wrong skills. Manager has no idea. Always use the full program keyword.',
    0.55, 4.45, 12.1, 0.5, size=9, color=AMBER)
box(sl, 0.4, 5.1, 12.5, 0.5, bg=AMBER_L, border=AMBER)
txt(sl, 'Adding a new alias — update BOTH files: api/pdf_generator.py (PROGRAM_ALIASES) AND public/index.html (PROGRAM_MAP). Missing one means local and production behave differently.',
    0.55, 5.15, 12.1, 0.4, size=9, color=AMBER)

# ── SLIDE 8 — FINAL SUMMARY ──────────────────────────────
sl = add_slide()
box(sl, 0, 0, 13.33, 7.5, bg=DARK)
txt(sl, 'Final Summary', 0.5, 0.4, 12.33, 0.7, size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

items = [
    'All 10 Rise Athletics gyms',
    'All 7 programs (Preschool, Junior A, Junior B, Advanced Junior, Level 1, Level 2, Level 3)',
    'All score types: blank (not yet), earned (star), mastery (gold star)',
    'Mixed scoring within a single class — different students, different patterns, same page',
    'Over-capacity multi-page classes (pg 1/2 and pg 2/2 splits)',
    'A/B/C class splits at the same time slot',
    'Real student data from actual iClassPro eval exports',
    '11,132+ individual data points verified cell-by-cell against source spreadsheet',
    'Zero errors found',
]
for i, item in enumerate(items):
    txt(sl, '✓  ' + item, 0.8, 1.3 + i * 0.48, 11.7, 0.44, size=12, color=GREEN_L)

banner_box(sl, 'The Star Chart Converter is verified and ready for production use across all 10 gyms.', 1.0, 6.55, 11.33, 0.7, bg=GREEN_L, fg=GREEN, size=14)

prs.save('stress tests/Star Chart Converter - Complete Audit Record.pptx')
print("Done.")
